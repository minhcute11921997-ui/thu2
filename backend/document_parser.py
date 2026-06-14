import csv
import json
import os
import re
import xml.etree.ElementTree as ET
from html.parser import HTMLParser
from pathlib import Path

import docx
import fitz

from config import CHUNK_OVERLAP, CHUNK_SIZE


class _HTMLTextExtractor(HTMLParser):
    def __init__(self):
        super().__init__()
        self.parts: list[str] = []

    def handle_data(self, data: str):
        text = data.strip()
        if text:
            self.parts.append(text)

    def get_text(self) -> str:
        return "\n".join(self.parts)


class DocumentParser:
    TEXT_EXTENSIONS = {
        ".txt",
        ".md",
        ".markdown",
        ".log",
        ".ini",
        ".conf",
        ".yaml",
        ".yml",
        ".sql",
        ".py",
        ".js",
        ".jsx",
        ".ts",
        ".tsx",
        ".java",
        ".cs",
        ".cpp",
        ".c",
        ".h",
        ".php",
        ".rb",
        ".go",
        ".rs",
        ".sh",
        ".bat",
        ".ps1",
    }
    SUPPORTED_EXTENSIONS = {
        ".pdf",
        ".docx",
        ".pptx",
        ".xlsx",
        ".xlsm",
        ".csv",
        ".tsv",
        ".json",
        ".xml",
        ".html",
        ".htm",
        ".rtf",
        *TEXT_EXTENSIONS,
    }

    @staticmethod
    def parse_pdf(file_path: str) -> str:
        text_parts = []
        ocr_error = None
        with fitz.open(file_path) as doc:
            for page in doc:
                page_text = page.get_text("text").strip()
                if page_text:
                    text_parts.append(page_text)

            if not text_parts:
                language = os.getenv("OCR_LANGUAGE", "vie+eng")
                for page in doc:
                    try:
                        textpage = page.get_textpage_ocr(language=language, full=True, dpi=200)
                        page_text = page.get_text("text", textpage=textpage).strip()
                        if page_text:
                            text_parts.append(page_text)
                    except RuntimeError as exc:
                        ocr_error = str(exc)
                        break

        if not text_parts and ocr_error:
            raise ValueError(
                "PDF nay khong co lop text. Can cai Tesseract OCR va bo ngon ngu vie/eng de doc PDF scan. "
                f"Chi tiet OCR: {ocr_error}"
            )

        return "\n".join(text_parts)

    @staticmethod
    def parse_docx(file_path: str) -> str:
        document = docx.Document(file_path)
        parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]

        for table in document.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))

        return "\n".join(parts)

    @staticmethod
    def parse_pptx(file_path: str) -> str:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise RuntimeError("Can cai python-pptx de doc file PPTX.") from exc

        presentation = Presentation(file_path)
        parts = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_parts.append(shape.text.strip())

                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            slide_parts.append(" | ".join(cells))

            if slide_parts:
                parts.append(f"Slide {slide_index}\n" + "\n".join(slide_parts))

        return "\n\n".join(parts)

    @staticmethod
    def parse_spreadsheet(file_path: str) -> str:
        try:
            from openpyxl import load_workbook
        except ImportError as exc:
            raise RuntimeError("Can cai openpyxl de doc file Excel.") from exc

        workbook = load_workbook(file_path, read_only=True, data_only=True)
        parts = []
        for sheet in workbook.worksheets:
            parts.append(f"Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                values = [str(value).strip() for value in row if value is not None and str(value).strip()]
                if values:
                    parts.append(" | ".join(values))

        workbook.close()
        return "\n".join(parts)

    @staticmethod
    def parse_text(file_path: str) -> str:
        encodings = ["utf-8-sig", "utf-8", "utf-16", "cp1258", "cp1252", "latin-1"]
        for encoding in encodings:
            try:
                with open(file_path, "r", encoding=encoding) as file:
                    return file.read()
            except UnicodeDecodeError:
                continue

        with open(file_path, "rb") as file:
            return file.read().decode("utf-8", errors="ignore")

    @classmethod
    def parse_delimited(cls, file_path: str, delimiter: str) -> str:
        text = cls.parse_text(file_path)
        rows = []
        for row in csv.reader(text.splitlines(), delimiter=delimiter):
            values = [value.strip() for value in row if value.strip()]
            if values:
                rows.append(" | ".join(values))

        return "\n".join(rows)

    @classmethod
    def parse_json(cls, file_path: str) -> str:
        text = cls.parse_text(file_path)
        data = json.loads(text)
        return json.dumps(data, ensure_ascii=False, indent=2)

    @classmethod
    def parse_xml(cls, file_path: str) -> str:
        text = cls.parse_text(file_path)
        root = ET.fromstring(text)
        parts = []
        for element in root.iter():
            value = (element.text or "").strip()
            if value:
                parts.append(f"{element.tag}: {value}")
        return "\n".join(parts)

    @classmethod
    def parse_html(cls, file_path: str) -> str:
        parser = _HTMLTextExtractor()
        parser.feed(cls.parse_text(file_path))
        return parser.get_text()

    @classmethod
    def parse_rtf(cls, file_path: str) -> str:
        text = cls.parse_text(file_path)
        text = re.sub(r"\\'[0-9a-fA-F]{2}", " ", text)
        text = re.sub(r"\\[a-zA-Z]+\d* ?", " ", text)
        text = text.replace("{", " ").replace("}", " ")
        return re.sub(r"\s+", " ", text).strip()

    @staticmethod
    def chunk_text(text: str) -> list[str]:
        words = text.replace("\n", " ").split()
        if not words:
            return []

        chunks = []
        step = max(1, CHUNK_SIZE - CHUNK_OVERLAP)
        for index in range(0, len(words), step):
            chunks.append(" ".join(words[index:index + CHUNK_SIZE]))

        return chunks

    @classmethod
    def supported_extensions(cls) -> list[str]:
        return sorted(cls.SUPPORTED_EXTENSIONS)

    @classmethod
    def process_file(cls, file_path: str) -> list[str]:
        ext = Path(file_path).suffix.lower()
        if ext == ".pdf":
            text = cls.parse_pdf(file_path)
        elif ext == ".docx":
            text = cls.parse_docx(file_path)
        elif ext == ".pptx":
            text = cls.parse_pptx(file_path)
        elif ext in [".xlsx", ".xlsm"]:
            text = cls.parse_spreadsheet(file_path)
        elif ext == ".csv":
            text = cls.parse_delimited(file_path, ",")
        elif ext == ".tsv":
            text = cls.parse_delimited(file_path, "\t")
        elif ext == ".json":
            text = cls.parse_json(file_path)
        elif ext == ".xml":
            text = cls.parse_xml(file_path)
        elif ext in [".html", ".htm"]:
            text = cls.parse_html(file_path)
        elif ext == ".rtf":
            text = cls.parse_rtf(file_path)
        elif ext in cls.TEXT_EXTENSIONS:
            text = cls.parse_text(file_path)
        else:
            raise ValueError(f"Dinh dang {ext or 'khong xac dinh'} khong duoc ho tro.")

        return cls.chunk_text(text)
